"""HTML(스크린샷 PNG) -> pptx 조립 스크립트.
사용법:
1. html/ 의 slide-*.html을 Playwright로 1920x1080 스크린샷 떠서 output/ 에 동일 이름의 .png로 저장
2. 이 스크립트를 실행하면:
   - output/kimkwanyoung_임시.pptx (개별 빌드본)
   - output/hajaCheck착수보고_최종.pptx (공용 PPTX와 병합된 최종본)
   을 파일명 번호(slide-NN)에 따라 동적으로 분석하여 생성합니다.
"""
import os
import glob
import re
import argparse
from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))
OUTPUT_DIR = os.path.join(HERE, "output")
DATA_DIR = os.path.join(HERE, "data")

INDIVIDUAL_PPTX_PATH = os.path.join(OUTPUT_DIR, "kimkwanyoung_임시.pptx")
FINAL_PPTX_PATH = os.path.join(OUTPUT_DIR, "hajaCheck착수보고_최종.pptx")
TEMPLATE_PPTX_PATH = os.path.join(DATA_DIR, "hajaCheck착수 보고.pptx")


def parse_slide_number(filename):
    """파일명(예: slide-05-plan.png)에서 슬라이드 번호(int)를 파싱합니다.
    파싱 실패 시 None을 반환합니다.
    """
    match = re.match(r"^slide-(\d+)", filename)
    if match:
        return int(match.group(1))
    return None


def get_screenshot_slides():
    """output/ 디렉터리 내의 slide-*.png 파일을 검색하여
    (슬라이드 번호, 파일 절대 경로, 파일명) 리스트를 정렬하여 반환합니다.
    """
    pattern = os.path.join(OUTPUT_DIR, "slide-*.png")
    files = glob.glob(pattern)
    
    slides = []
    for path in files:
        name = os.path.basename(path)
        num = parse_slide_number(name)
        if num is not None:
            slides.append((num, path, name))
            
    # 슬라이드 번호 기준 오름차순 정렬
    slides.sort(key=lambda x: x[0])
    return slides


def clear_slide_shapes(slide):
    """슬라이드 내의 모든 도형 및 텍스트 박스를 안전하게 삭제합니다."""
    shapes = list(slide.shapes)
    for shape in reversed(shapes):
        el = shape.element
        el.getparent().remove(el)


def build_individual():
    """캡처된 모든 슬라이드만 모아 임시 PPTX를 빌드합니다."""
    slides = get_screenshot_slides()
    if not slides:
        print("경고: 조립할 스크린샷 이미지(slide-*.png)가 output/ 에 없습니다.")
        return False

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    print("임시 개별 PPTX 조립 중...")
    for num, path, name in slides:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        print(f"  - 슬라이드 {num} 추가 완료 ({name})")

    prs.save(INDIVIDUAL_PPTX_PATH)
    print(f"개별 PPTX 저장 완료: {INDIVIDUAL_PPTX_PATH}")
    return True


def build_merged():
    """공용 PPTX 원본을 불러와 파일명 번호에 해당하는 슬라이드를 덮어쓰고,
    기존 PPTX 장수를 초과하는 슬라이드는 뒤에 덧붙여서 최종본을 만듭니다.
    """
    if not os.path.exists(TEMPLATE_PPTX_PATH):
        print(f"경고: 공용 템플릿 PPTX 파일이 존재하지 않습니다: {TEMPLATE_PPTX_PATH}")
        print("병합 빌드를 건너뜁니다.")
        return False

    slides = get_screenshot_slides()
    if not slides:
        print("경고: 병합할 스크린샷 이미지(slide-*.png)가 output/ 에 없습니다.")
        return False

    print("공용 PPTX를 베이스로 최종 병합 PPTX 빌드 중...")
    prs = Presentation(TEMPLATE_PPTX_PATH)
    width = prs.slide_width
    height = prs.slide_height
    blank_layout = prs.slide_layouts[6]

    # 공용 PPTX의 기본 슬라이드 개수 (예: 8)
    original_slide_count = len(prs.slides)
    print(f"  - 공용 템플릿 슬라이드 수: {original_slide_count}")

    # 슬라이드를 교체(Overwrite)할 그룹과 새로 추가(Append)할 그룹으로 분류
    to_replace = []
    to_append = []

    for num, path, name in slides:
        if num <= original_slide_count:
            to_replace.append((num, path, name))
        else:
            to_append.append((num, path, name))

    # 1. 기존 슬라이드 교체
    for num, path, name in to_replace:
        idx = num - 1  # 0-indexed 변환
        slide = prs.slides[idx]
        clear_slide_shapes(slide)
        slide.shapes.add_picture(path, 0, 0, width=width, height=height)
        print(f"  - 슬라이드 {num} 교체 완료: {name}")

    # 2. 신규 슬라이드 추가
    for num, path, name in to_append:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(path, 0, 0, width=width, height=height)
        print(f"  - 신규 슬라이드 {num} 추가 완료: {name}")

    prs.save(FINAL_PPTX_PATH)
    print(f"최종 병합 PPTX 저장 완료: {FINAL_PPTX_PATH}")
    return True


def main():
    parser = argparse.ArgumentParser(description="동적 파일명 기반 PPTX 조립 및 병합 스크립트")
    parser.add_argument("--individual-only", action="store_true", help="개별 임시 PPTX만 빌드합니다.")
    args = parser.parse_args()

    if not os.path.exists(OUTPUT_DIR):
        os.makedirs(OUTPUT_DIR)

    # 개별 빌드 수행
    success_ind = build_individual()
    
    # 병합 빌드 수행 (옵션 지정 안 했을 때만)
    if not args.individual_only and success_ind:
        build_merged()


if __name__ == "__main__":
    main()
